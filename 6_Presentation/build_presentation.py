"""
Build the final presentation for CSC 3221 — Mobile Money Transaction Analysis.
12 slides with demographic integration, professional styling, and speaker notes.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Color palette ---
DARK_BLUE   = RGBColor(0x1B, 0x2A, 0x4A)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BLUE  = RGBColor(0xD6, 0xEA, 0xF8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)
MED_GRAY    = RGBColor(0x66, 0x66, 0x66)
ORANGE      = RGBColor(0xE6, 0x7E, 0x22)
GREEN       = RGBColor(0x27, 0xAE, 0x60)
RED_ACCENT  = RGBColor(0xC0, 0x39, 0x2B)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
TOTAL_SLIDES = 12


def add_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=18,
                    color=DARK_GRAY, bullet_color=ACCENT_BLUE, spacing=Pt(8),
                    font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = spacing
        p.level = 0
        from pptx.oxml.ns import qn
        from lxml import etree
        pPr = p._p.get_or_add_pPr()
        for bn in pPr.findall(qn('a:buNone')):
            pPr.remove(bn)
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', '●')
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        srgb = etree.SubElement(buClr, qn('a:srgbClr'))
        srgb.set('val', str(bullet_color))
        buSzPct = etree.SubElement(pPr, qn('a:buSzPct'))
        buSzPct.set('val', '70000')
    return txBox


def add_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = text


def add_section_number(slide, number):
    add_textbox(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4),
                f'{number}/{TOTAL_SLIDES}', font_size=11, color=MED_GRAY,
                alignment=PP_ALIGN.RIGHT)


def add_accent_bar(slide, top=Inches(1.55), width=Inches(1.5)):
    add_shape_rect(slide, Inches(0.8), top, width, Inches(0.06), ACCENT_BLUE)


def add_slide_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(1.0),
                title, font_size=32, color=DARK_BLUE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(1.15), Inches(11.5), Inches(0.5),
                    subtitle, font_size=16, color=MED_GRAY)
    add_accent_bar(slide)


# ===== BUILD PRESENTATION =====
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT
blank_layout = prs.slide_layouts[6]


# ========== SLIDE 1: TITLE ==========
slide = prs.slides.add_slide(blank_layout)
add_background(slide, DARK_BLUE)
add_shape_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15), ACCENT_BLUE)
add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2),
            'Mobile Money Transaction Analysis', font_size=44, color=WHITE, bold=True)
add_textbox(slide, Inches(1.5), Inches(3.0), Inches(10.3), Inches(0.8),
            'User Activity Classification in Cameroon', font_size=28, color=LIGHT_BLUE)
add_shape_rect(slide, Inches(1.5), Inches(3.9), Inches(3.0), Inches(0.04), ACCENT_BLUE)
add_textbox(slide, Inches(1.5), Inches(4.3), Inches(10.3), Inches(0.5),
            'CSC 3221 — Introduction to Data Science', font_size=18, color=LIGHT_BLUE)
add_textbox(slide, Inches(1.5), Inches(4.9), Inches(10.3), Inches(0.5),
            'Final Assessment  |  May 2026', font_size=16, color=MED_GRAY)
add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5),
            'Dr. Fotsing Kuetche — ICT University', font_size=14, color=MED_GRAY)
add_shape_rect(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), ACCENT_BLUE)
add_speaker_notes(slide,
    "Welcome everyone. Today we present our analysis of mobile money transaction "
    "patterns in Cameroon. Our goal was to classify users into behavioral segments "
    "using both transaction behavior and demographic attributes. "
    "This has direct applications in customer segmentation, credit risk assessment, "
    "and financial inclusion strategy.")


# ========== SLIDE 2: PROBLEM STATEMENT ==========
slide = prs.slides.add_slide(blank_layout)
add_background(slide, WHITE)
add_slide_title(slide, 'Problem Statement', 'What are we trying to solve?')
add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(7.0), Inches(4.5), [
    'Goal: Classify mobile money users into Low / Medium / High activity',
    'Multi-class classification using behavioral + demographic features',
    'Target defined by quantile thresholds on transaction count',
    'Must outperform naive baselines with no data leakage',
    'New: Assess impact of demographic data on classification accuracy',
], font_size=20)

# Right-side info cards
for i, (label, val, clr) in enumerate([
    ('Classes', '3', ACCENT_BLUE),
    ('Users', '69', ORANGE),
    ('Features', '29', GREEN),
]):
    top = Inches(2.2 + i * 1.4)
    card = add_shape_rect(slide, Inches(9.0), top, Inches(3.5), Inches(1.1), clr)
    add_textbox(slide, Inches(9.2), top + Inches(0.1), Inches(3.0), Inches(0.4),
                label, font_size=14, color=WHITE)
    add_textbox(slide, Inches(9.2), top + Inches(0.45), Inches(3.0), Inches(0.5),
                val, font_size=36, color=WHITE, bold=True)

add_section_number(slide, 2)
add_speaker_notes(slide,
    "We defined three activity classes using quantile-based thresholds: Low is 55 or fewer "
    "transactions, Medium is 56 to 204, and High is above 204. The classification uses "
    "15 behavioral features describing how users transact, plus 14 demographic features "
    "describing who they are. The key challenge: only 69 users with a very small dataset.")


# ========== SLIDE 3: DATASET OVERVIEW ==========
slide = prs.slides.add_slide(blank_layout)
add_background(slide, WHITE)
add_slide_title(slide, 'Dataset Overview', 'Two data sources integrated')
add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), [
    '20,530 transactions from 69 users (51 original + 18 anonymized)',
    'Date range: 2023–2026 | MobileMoney (96%) + OrangeMoney (4%)',
    'Demographic questionnaire: 10 users (extended to 69 via simulation)',
    'Fields: age, gender, profession, education, income, geography',
    'All records tagged: real (10) vs simulated (59)',
], font_size=18)

# Stats grid
for j, (lbl, val) in enumerate([
    ('Transactions', '20,530'), ('Users', '69'),
    ('Behavioral Feat.', '15'), ('Demographic Feat.', '14'),
]):
    col = j % 2
    row = j // 2
    left = Inches(7.5 + col * 2.7)
    top = Inches(2.2 + row * 1.8)
    add_shape_rect(slide, left, top, Inches(2.4), Inches(1.4), LIGHT_BLUE)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.15), Inches(2.1), Inches(0.4),
                lbl, font_size=12, color=ACCENT_BLUE)
    add_textbox(slide, left + Inches(0.15), top + Inches(0.5), Inches(2.1), Inches(0.6),
                val, font_size=28, color=DARK_BLUE, bold=True)

add_section_number(slide, 3)
add_speaker_notes(slide,
    "Our behavioral data comes from SMS transaction messages exported by 69 users. "
    "We augmented this with a demographic questionnaire completed by 10 users. "
    "For the remaining 59, we simulated plausible demographics based on the 10 real "
    "responses and Cameroon population statistics. This is clearly documented as simulated "
    "data. The behavioral features describe HOW users transact, while demographic features "
    "describe WHO they are.")


# ========== SLIDE 4: DATA CLEANING SUMMARY ==========
slide = prs.slides.add_slide(blank_layout)
add_background(slide, WHITE)
add_slide_title(slide, 'Data Cleaning Summary', 'Behavioral + Demographic pipelines')

# Left column: Behavioral
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.4),
            'Behavioral Pipeline', font_size=20, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(3.0), [
    'Parsed SMS messages into structured transactions',
    'Standardized 6 transaction types + IN/OUT direction',
    'Engineered 15 user-level features (amounts, timing, ratios)',
    'Excluded leakage features (frequency, velocity, count)',
], font_size=16)

# Right column: Demographic
add_textbox(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.4),
            'Demographic Pipeline', font_size=20, color=ORANGE, bold=True)
add_bullet_list(slide, Inches(7.0), Inches(2.5), Inches(5.5), Inches(3.0), [
    'Normalized age ranges, gender, income brackets',
    'Grouped 8 professions into 5 standard categories',
    'Standardized education to 4 levels',
    'Fixed encoding issues (Yaoundé), simplified geography',
], font_size=16, bullet_color=ORANGE)

# Leakage box
add_shape_rect(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.2), LIGHT_BLUE)
add_textbox(slide, Inches(1.0), Inches(5.6), Inches(11.0), Inches(0.4),
            'Data Leakage Prevention', font_size=16, color=DARK_BLUE, bold=True)
add_textbox(slide, Inches(1.0), Inches(6.0), Inches(11.0), Inches(0.5),
            'Excluded: txn_count, monthly_txn_frequency, txn_velocity_7d, data_source (all overlap with target)',
            font_size=14, color=DARK_GRAY)

add_section_number(slide, 4)
add_speaker_notes(slide,
    "Two parallel cleaning pipelines. The behavioral pipeline parses raw SMS messages "
    "into structured transactions and aggregates to user-level features. The demographic "
    "pipeline standardizes questionnaire responses. Critically, we excluded all frequency-based "
    "features to prevent data leakage, since our target is based on transaction count.")


# ========== SLIDE 5: EDA HIGHLIGHTS ==========
slide = prs.slides.add_slide(blank_layout)
add_background(slide, WHITE)
add_slide_title(slide, 'EDA Highlights', 'Key patterns in the data')

add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5), [
    'Micro-transactions dominate: median = 500 XAF (~$0.80), mean = 7,454 XAF (skewness > 12)',
    '77.5% outbound transactions (transfers, airtime, payments)',
    'Evening peak (17:00–20:00) consistent across all days',
    'High-activity users: many small transactions; Low-activity: fewer but larger',
    'Balanced classes: ~23 users per activity level (quantile-based split)',
], font_size=18)

# Key insight box
add_shape_rect(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.2), LIGHT_BLUE)
add_textbox(slide, Inches(1.0), Inches(5.65), Inches(11.0), Inches(0.9),
            'Key insight: Transaction frequency is more stable than amount for classification. '
            'How a user transacts (types, timing, direction) predicts engagement level.',
            font_size=16, color=DARK_BLUE)

add_section_number(slide, 5)
add_speaker_notes(slide,
    "Our EDA revealed that mobile money in Cameroon is dominated by small, frequent "
    "transactions. The median is just 500 XAF. Most transactions are outbound, with "
    "transfers and airtime topping the list. We identified clear behavioral differences "
    "between activity classes that the model can exploit.")


# ========== SLIDE 6: MODELING APPROACH ==========
slide = prs.slides.add_slide(blank_layout)
add_background(slide, WHITE)
add_slide_title(slide, 'Modeling Approach', '3 models, 2 feature sets, stratified evaluation')

# Model cards
for i, (name, desc, clr) in enumerate([
    ('Logistic Regression', 'Linear baseline\nC=1, balanced weights\nFull interpretability', ACCENT_BLUE),
    ('Random Forest', 'Ensemble of 200 trees\nmax_depth=5\nFeature importance', GREEN),
    ('XGBoost', 'Gradient boosting\nmax_depth=3, lr=0.1\nBest prior performance', ORANGE),
]):
    left = Inches(0.8 + i * 4.0)
    card = add_shape_rect(slide, left, Inches(2.2), Inches(3.6), Inches(2.0), clr)
    add_textbox(slide, left + Inches(0.2), Inches(2.35), Inches(3.2), Inches(0.4),
                name, font_size=18, color=WHITE, bold=True)
    add_textbox(slide, left + Inches(0.2), Inches(2.85), Inches(3.2), Inches(1.2),
                desc, font_size=13, color=WHITE)

add_bullet_list(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(2.5), [
    '70/30 stratified split | 5-fold stratified cross-validation',
    'Comparison: Behavioral Only (15 features) vs Behavioral + Demographic (29 features)',
    'Primary metric: Macro F1 (treats all 3 classes equally)',
    'random_state=42 throughout for reproducibility',
], font_size=16)

add_section_number(slide, 6)
add_speaker_notes(slide,
    "We trained three models: Logistic Regression as an interpretable baseline, "
    "Random Forest as a robust ensemble, and XGBoost as the highest-performing "
    "gradient boosting approach. Each was evaluated with and without demographic "
    "features to measure the impact of socioeconomic data.")


# ========== SLIDE 7: RESULTS ==========
slide = prs.slides.add_slide(blank_layout)
add_background(slide, WHITE)
add_slide_title(slide, 'Results', 'Test set performance (F1 Macro)')

# Results table
headers = ['Model', 'Behav F1m', 'Combined F1m', 'Delta']
data = [
    ['Majority Baseline', '0.167', '—', '—'],
    ['Stratified Baseline', '0.315', '—', '—'],
    ['Logistic Regression', '0.430', '0.423', '-0.007'],
    ['Random Forest', '0.418', '0.470', '+0.052'],
    ['XGBoost', '0.570', '0.529', '-0.041'],
]

col_widths = [Inches(3.5), Inches(2.2), Inches(2.5), Inches(2.0)]
table_left = Inches(1.2)
table_top = Inches(2.2)
row_h = Inches(0.5)

# Header row
for j, (hdr, cw) in enumerate(zip(headers, col_widths)):
    x = table_left + sum(cw.inches for cw in col_widths[:j]) * 914400
    x = table_left + Inches(sum(c.inches for c in col_widths[:j]))
    add_shape_rect(slide, x, table_top, cw, row_h, DARK_BLUE)
    add_textbox(slide, x + Inches(0.1), table_top + Inches(0.05), cw - Inches(0.2), row_h,
                hdr, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Data rows
for ri, row in enumerate(data):
    y = table_top + row_h * (ri + 1)
    bg = LIGHT_BLUE if ri % 2 == 0 else WHITE
    is_best = ri == 4
    for j, (val, cw) in enumerate(zip(row, col_widths)):
        x = table_left + Inches(sum(c.inches for c in col_widths[:j]))
        add_shape_rect(slide, x, y, cw, row_h, bg)
        fc = DARK_BLUE if is_best else DARK_GRAY
        bld = is_best
        add_textbox(slide, x + Inches(0.1), y + Inches(0.05), cw - Inches(0.2), row_h,
                    val, font_size=14, color=fc, bold=bld, alignment=PP_ALIGN.CENTER)

# Key takeaway
add_shape_rect(slide, Inches(1.2), Inches(5.4), Inches(10.2), Inches(1.2), LIGHT_BLUE)
add_textbox(slide, Inches(1.4), Inches(5.5), Inches(9.8), Inches(1.0),
            'XGBoost (Behavioral Only) leads at F1m = 0.570. '
            'Demographics improved RF (+0.05) but slightly decreased XGBoost (-0.04). '
            'All trained models exceed both baselines.',
            font_size=16, color=DARK_BLUE)

add_section_number(slide, 7)
add_speaker_notes(slide,
    "XGBoost with behavioral features alone achieved the highest F1 macro of 0.57. "
    "When we added demographics, Random Forest improved but XGBoost slightly decreased. "
    "This is because 59 of 69 demographics are simulated, adding noise. With real "
    "demographic data, we would expect more consistent improvement. All models beat baselines.")


# ========== SLIDE 8: FEATURE IMPORTANCE ==========
slide = prs.slides.add_slide(blank_layout)
add_background(slide, WHITE)
add_slide_title(slide, 'Feature Importance', 'What drives the predictions?')

# Behavioral features
add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.4),
            'Top Behavioral Features', font_size=18, color=ACCENT_BLUE, bold=True)
add_bullet_list(slide, Inches(0.8), Inches(2.5), Inches(5.5), Inches(3.5), [
    'pct_depot (0.104) — deposit frequency distinguishes levels',
    'avg_sr_ratio (0.052) — send/receive balance',
    'pct_transfert (0.055) — transfer proportion',
    'mean_amount (0.051) — transaction size profile',
    'avg_hour (0.079 RF) — usage timing patterns',
], font_size=16)

# Demographic features
add_textbox(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.4),
            'Top Demographic Features', font_size=18, color=ORANGE, bold=True)
add_bullet_list(slide, Inches(7.0), Inches(2.5), Inches(5.5), Inches(3.5), [
    'prof_Unemployed (0.046) — employment status',
    'prof_Self-Employed (0.045) — business users',
    'education_ordinal (0.044) — education level',
    'geo_Urban (0.040) — location type',
    'household_size — family network size',
], font_size=16, bullet_color=ORANGE)

add_shape_rect(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.8), LIGHT_BLUE)
add_textbox(slide, Inches(1.0), Inches(6.1), Inches(11.0), Inches(0.6),
            'Behavioral features dominate the top ranks. Demographic features provide secondary signal.',
            font_size=15, color=DARK_BLUE)

add_section_number(slide, 8)
add_speaker_notes(slide,
    "Feature importance analysis reveals that behavioral features are the primary "
    "predictors. Deposit frequency, send-receive ratio, and transfer patterns are "
    "the top signals. Among demographics, profession and education show the strongest "
    "signal, suggesting that occupation type influences mobile money behavior.")


# ========== SLIDE 9: KEY INSIGHTS ==========
slide = prs.slides.add_slide(blank_layout)
add_background(slide, DARK_BLUE)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
            'Key Insights', font_size=32, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.06), ACCENT_BLUE)

# Three user profiles
profiles = [
    ('HIGH Activity', '> 204 transactions',
     'Daily power users\nSmall frequent transactions\nHigh send/receive ratio\nAirtime + transfers\nLarger households',
     GREEN),
    ('MEDIUM Activity', '56–204 transactions',
     'Regular but not daily\nModerate amounts\nMixed transaction types\nTransitional group\nHardest to classify (F1=0.40)',
     ORANGE),
    ('LOW Activity', '≤ 55 transactions',
     'Infrequent, purpose-driven\nFewer but larger transactions\nHigh balance, lower ratio\nHigher education\nSpecific-use pattern',
     RED_ACCENT),
]
for i, (title, sub, desc, clr) in enumerate(profiles):
    left = Inches(0.8 + i * 4.0)
    card = add_shape_rect(slide, left, Inches(1.8), Inches(3.6), Inches(4.5), clr)
    add_textbox(slide, left + Inches(0.2), Inches(1.95), Inches(3.2), Inches(0.4),
                title, font_size=20, color=WHITE, bold=True)
    add_textbox(slide, left + Inches(0.2), Inches(2.35), Inches(3.2), Inches(0.35),
                sub, font_size=13, color=WHITE)
    add_textbox(slide, left + Inches(0.2), Inches(2.85), Inches(3.2), Inches(3.0),
                desc, font_size=14, color=WHITE)

add_textbox(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.5),
            'How a user transacts is more predictive than who they are.',
            font_size=18, color=LIGHT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_section_number(slide, 9)
add_speaker_notes(slide,
    "We identified three distinct user profiles. High-activity users are daily power users "
    "making many small transactions. Low-activity users are infrequent but purposeful, with "
    "larger transaction amounts. Medium users are the hardest to classify because they blend "
    "characteristics of both groups. The key takeaway: behavioral features are more "
    "predictive than demographics for this task.")


# ========== SLIDE 10: BUSINESS RECOMMENDATIONS ==========
slide = prs.slides.add_slide(blank_layout)
add_background(slide, WHITE)
add_slide_title(slide, 'Business Recommendations', 'From classification to action')

# Application cards
apps = [
    ('Customer Segmentation', 'Classify new users within first\nmonth. Tailor products per tier:\npremium for High, re-engagement\nfor Low.', ACCENT_BLUE),
    ('Credit Risk Scoring', 'High-activity users show consistent\nengagement — positive signal for\nmicro-lending. Behavioral features\ncomplement traditional scores.', GREEN),
    ('Targeted Marketing', 'Urban self-employed High users:\nbusiness accounts. Students:\nairtime bundles. Educated Low:\nsavings products.', ORANGE),
    ('Churn Detection', 'Monitor activity class transitions.\nDeclining velocity signals churn.\nMisclassified users may indicate\nanomalous behavior.', RED_ACCENT),
]
for i, (title, desc, clr) in enumerate(apps):
    col = i % 2
    row = i // 2
    left = Inches(0.8 + col * 6.2)
    top = Inches(2.0 + row * 2.5)
    add_shape_rect(slide, left, top, Inches(5.8), Inches(2.1), clr)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.1), Inches(5.4), Inches(0.4),
                title, font_size=18, color=WHITE, bold=True)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.55), Inches(5.4), Inches(1.4),
                desc, font_size=14, color=WHITE)

add_section_number(slide, 10)
add_speaker_notes(slide,
    "Four business applications. First, customer segmentation: classify users early and "
    "tailor products. Second, credit risk: high-activity users are better lending candidates. "
    "Third, targeted marketing using the demographic-behavioral intersection. Fourth, "
    "churn detection by monitoring activity class transitions over time.")


# ========== SLIDE 11: LIMITATIONS ==========
slide = prs.slides.add_slide(blank_layout)
add_background(slide, WHITE)
add_slide_title(slide, 'Limitations', 'What constrains these results?')

add_bullet_list(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5), [
    'Small dataset: only 69 users — models have limited learning capacity',
    'Simulated demographics: 59/69 are synthetic — dampens real demographic signal',
    'Feature dimensionality: 29 features on 69 users pushes model capacity',
    'No temporal split: random split doesn’t simulate real deployment',
    'Leakage tradeoff: excluding frequency features limits maximum performance',
    'Medium class overlap: boundary users are inherently hard to classify',
], font_size=18)

add_shape_rect(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.0), LIGHT_BLUE)
add_textbox(slide, Inches(1.0), Inches(5.9), Inches(11.0), Inches(0.8),
            'Primary recommendation: Expand to 200+ users with real demographics to stabilize '
            'model performance and enable genuine demographic analysis.',
            font_size=15, color=DARK_BLUE)

add_section_number(slide, 11)
add_speaker_notes(slide,
    "The main limitation is dataset size. With 69 users, models are inherently unstable. "
    "The simulated demographics add features but may introduce noise. We recommend "
    "expanding to 200+ users with complete real demographic data. A temporal split "
    "approach would also improve reliability by simulating real-world deployment.")


# ========== SLIDE 12: CONCLUSION ==========
slide = prs.slides.add_slide(blank_layout)
add_background(slide, DARK_BLUE)
add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.8),
            'Conclusion', font_size=32, color=WHITE, bold=True)
add_shape_rect(slide, Inches(0.8), Inches(1.3), Inches(1.5), Inches(0.06), ACCENT_BLUE)

# Summary points
findings = [
    'Behavioral features carry genuine predictive signal for user classification',
    'XGBoost: Test F1 (macro) = 0.570, well above baselines (0.17 / 0.32)',
    'Demographics: mixed impact (+0.05 RF, -0.04 XGBoost) due to 85% simulated data',
    'Top predictors: deposit %, send/receive ratio, transfer %, transaction amount',
    'Three distinct user profiles identified: power users, regular users, occasional users',
]
add_bullet_list(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(3.5), findings,
                font_size=20, color=WHITE, bullet_color=ACCENT_BLUE)

# Next steps box
add_shape_rect(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5), ACCENT_BLUE)
add_textbox(slide, Inches(1.0), Inches(5.3), Inches(11.0), Inches(0.4),
            'Next Steps', font_size=18, color=WHITE, bold=True)
add_textbox(slide, Inches(1.0), Inches(5.7), Inches(11.0), Inches(0.8),
            '1. Collect real demographics for all users  |  '
            '2. Expand to 200+ users  |  '
            '3. Implement temporal split  |  '
            '4. Deploy batch scoring pipeline',
            font_size=14, color=WHITE)

add_shape_rect(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.15), ACCENT_BLUE)
add_speaker_notes(slide,
    "In conclusion, we demonstrated that mobile money transaction behavior alone carries "
    "genuine predictive signal for user activity classification. XGBoost achieved an F1 "
    "of 0.57, well above baselines. Demographic integration showed mixed results due to "
    "mostly simulated data, but profession and education provided secondary signals. "
    "The three user profiles we identified have direct business applications. Thank you.")


# ===== SAVE =====
output_path = os.path.join(os.path.dirname(__file__), 'Presentation_Slides.pptx')
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
